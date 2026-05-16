from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
OUTPUT_DIR = ROOT / "outputs" / "Chennai_Restaurant_答辩PPT_Editable"
PPT_PATH = OUTPUT_DIR / "Chennai_Restaurant_答辩PPT_Editable.pptx"

FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"

BLUE = RGBColor(9, 52, 133)
BLUE_2 = RGBColor(29, 88, 179)
PALE_BG = RGBColor(246, 249, 255)
LIGHT_BLUE = RGBColor(233, 241, 255)
LINE = RGBColor(214, 225, 242)
TEXT = RGBColor(29, 36, 51)
MUTED = RGBColor(92, 102, 122)
RED = RGBColor(192, 0, 0)
GOLD = RGBColor(217, 178, 94)
WHITE = RGBColor(255, 255, 255)
CARD = RGBColor(255, 255, 255)
SOFT_YELLOW = RGBColor(255, 248, 233)
SOFT_GREEN = RGBColor(239, 247, 239)
SOFT_RED = RGBColor(255, 240, 240)


def cm(value: float):
    return Cm(value)


def ensure_output() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = PALE_BG


def add_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, cm(x), cm(y), cm(w), cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.5)
    return shp


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=20,
    bold=False,
    color=TEXT,
    font=FONT_CN,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.12,
):
    tb = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = cm(margin)
    tf.margin_right = cm(margin)
    tf.margin_top = cm(margin)
    tf.margin_bottom = cm(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(
    slide,
    x,
    y,
    w,
    h,
    items: Iterable[str],
    size=18,
    color=TEXT,
    bullet_color=RED,
    spacing=1.18,
):
    tb = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = cm(0.15)
    tf.margin_right = cm(0.1)
    tf.margin_top = cm(0.05)
    tf.margin_bottom = cm(0.05)
    tf.clear()
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.name = FONT_CN
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.line_spacing = spacing
        try:
            p._pPr.insert(0, p._element._new_buChar(char="•"))
        except Exception:
            pass
    return tb


def style_header(slide, section_no: str, title: str, subtitle: str, page_no: int) -> None:
    set_slide_bg(slide)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cm(0), cm(0), cm(33.867), cm(0.18))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    add_textbox(slide, 0.75, 0.35, 1.2, 0.9, section_no, size=26, bold=True, color=BLUE, font=FONT_EN)
    add_textbox(slide, 2.1, 0.3, 18, 1.0, title, size=24, bold=True, color=BLUE)
    add_textbox(slide, 0.9, 1.3, 20, 0.7, subtitle, size=11, color=MUTED)
    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cm(0), cm(18.55), cm(33.867), cm(0.5))
    footer.fill.solid()
    footer.fill.fore_color.rgb = BLUE
    footer.line.fill.background()
    add_textbox(slide, 31.6, 18.54, 1.3, 0.4, str(page_no).zfill(2), size=12, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def add_chip(slide, x, y, w, h, text, fill=BLUE, color=WHITE, size=10, bold=True):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, cm(x), cm(y), cm(w), cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    add_textbox(slide, x + 0.12, y + 0.03, w - 0.24, h - 0.06, text, size=size, bold=bold, color=color, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return shp


def add_kpi(slide, x, y, w, h, title, value, note):
    add_rect(slide, x, y, w, h, fill=SOFT_YELLOW, line=RGBColor(228, 212, 174), radius=True)
    add_textbox(slide, x + 0.25, y + 0.18, w - 0.5, 0.45, title, size=10, color=MUTED, bold=True)
    add_textbox(slide, x + 0.25, y + 0.75, w - 0.5, 0.85, value, size=22, bold=True, color=TEXT, font=FONT_EN)
    add_textbox(slide, x + 0.25, y + 1.6, w - 0.5, h - 1.75, note, size=9, color=MUTED)


def add_panel(slide, x, y, w, h, title=None, fill=CARD):
    add_rect(slide, x, y, w, h, fill=fill, line=LINE, radius=True)
    if title:
        title_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, cm(x), cm(y), cm(w), cm(0.8))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = LIGHT_BLUE
        title_bar.line.fill.background()
        add_textbox(slide, x + 0.2, y + 0.12, w - 0.4, 0.45, title, size=12, bold=True, color=BLUE)


def add_picture(slide, image_path: Path, x, y, w, h):
    slide.shapes.add_picture(str(image_path), cm(x), cm(y), cm(w), cm(h))


def add_callout(slide, x, y, w, h, title, body, fill=SOFT_GREEN):
    add_rect(slide, x, y, w, h, fill=fill, line=LINE, radius=True)
    add_textbox(slide, x + 0.2, y + 0.15, w - 0.4, 0.45, title, size=11, bold=True, color=BLUE_2)
    add_textbox(slide, x + 0.2, y + 0.65, w - 0.4, h - 0.8, body, size=10, color=TEXT)


def write_notes(slide, notes: str):
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = notes


def cover_slide(prs, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cm(0), cm(0), cm(33.867), cm(19.05))
    hero.fill.solid()
    hero.fill.fore_color.rgb = WHITE
    hero.line.fill.background()
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cm(0), cm(2.2), cm(33.867), cm(0.22))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    add_textbox(slide, 1.0, 3.0, 14.5, 1.3, "Chennai Restaurant", size=28, bold=True, color=BLUE, font=FONT_EN)
    add_textbox(slide, 1.0, 4.4, 16.5, 1.4, "数据分析与可视化平台", size=30, bold=True, color=TEXT)
    add_textbox(slide, 1.1, 6.0, 16.0, 1.2, "基于 Zomato Chennai 餐饮数据的空间洞察、市场细分与评分预测研究", size=14, color=MUTED)
    add_textbox(slide, 1.1, 7.5, 14.8, 2.0, "从研究背景、数据准备、系统实现到 EDA 与机器学习解释，形成一套可复用的城市餐饮分析范式。", size=16, bold=True, color=TEXT)
    add_chip(slide, 1.1, 10.0, 4.4, 0.85, "01 项目背景与意义")
    add_chip(slide, 5.9, 10.0, 4.2, 0.85, "02 数据集与预处理")
    add_chip(slide, 10.4, 10.0, 4.2, 0.85, "03 技术栈与架构")
    add_chip(slide, 14.9, 10.0, 4.5, 0.85, "04 EDA 与评分预测")
    add_panel(slide, 20.2, 2.8, 12.6, 7.0, "平台主界面")
    add_picture(slide, ASSETS / "02 数据集" / "01 数据集截图.png", 20.6, 3.7, 11.8, 5.3)
    add_panel(slide, 20.2, 10.4, 12.6, 6.5, "核心证据预览")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "overview" / "overview.png", 20.55, 11.25, 5.7, 4.9)
    add_picture(slide, ASSETS / "04 EDA数据分析" / "predict" / "模型排行榜.png", 26.5, 11.25, 5.8, 4.9)
    add_textbox(slide, 1.1, 17.8, 5.0, 0.5, "毕业论文答辩", size=11, bold=True, color=WHITE)
    add_textbox(slide, 31.5, 17.8, 1.2, 0.5, str(page_no).zfill(2), size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cm(0), cm(18.55), cm(33.867), cm(0.5))
    footer.fill.solid()
    footer.fill.fore_color.rgb = BLUE
    footer.line.fill.background()
    write_notes(slide, "封面页先介绍研究主题、对象和整场答辩的四个章节。这里建议先强调本项目不仅是可视化展示，更包含空间分析和评分预测解释。")


def agenda_slide(prs, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "00", "汇报目录", "将答辩内容拆为五个层次，避免图表堆叠导致信息拥堵。", page_no)
    items = [
        ("01", "项目背景与研究意义", "研究问题、业务价值、项目目标"),
        ("02", "数据集与预处理", "数据来源、字段结构、长格式拆分、特征工程"),
        ("03", "技术栈与平台实现", "Flask 架构、前后端技术、分析服务组织"),
        ("04", "探索性数据分析", "整体概况、空间差异、业态差异"),
        ("05", "评分预测与总结", "模型比较、误差解释、主要结论与展望"),
    ]
    y = 3.0
    for idx, (num, title, desc) in enumerate(items):
        add_panel(slide, 1.2, y, 31.2, 2.45)
        add_chip(slide, 1.65, y + 0.55, 1.55, 0.85, num)
        add_textbox(slide, 3.6, y + 0.42, 8.4, 0.55, title, size=18, bold=True, color=TEXT)
        add_textbox(slide, 3.6, y + 1.1, 19.0, 0.7, desc, size=11, color=MUTED)
        if idx == 3:
            add_textbox(slide, 24.8, y + 0.75, 6.6, 0.8, "本次会拆成多页展示，确保每张图能看清。", size=10, color=RED, bold=True, align=PP_ALIGN.RIGHT)
        y += 2.8
    write_notes(slide, "目录页要明确说明这次重做后的结构更细，特别是 EDA 和模型预测部分会拆页展开，保证图表和文字不会挤在一张图里。")


def background_slides(prs, start_page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "01", "研究背景", "先交代为什么选择 Chennai 餐饮市场作为研究对象。", start_page)
    add_panel(slide, 1.0, 2.3, 15.2, 7.2, "研究背景与问题来源")
    add_bullets(slide, 1.3, 3.2, 14.4, 5.6, [
        "随着互联网平台积累了大量餐厅评分、菜系、地理位置与服务信息，城市餐饮市场具备了数据化分析基础。",
        "Chennai 作为南印度重要都市，人口密集、文化多元、餐饮供给丰富，适合进行城市级餐饮结构研究。",
        "消费者在选店时常面对信息过载问题，商家在选址与业态判断时也缺少结构化证据支持。",
        "因此，本项目尝试将平台数据转化为“区域差异 + 业态差异 + 评分驱动”的综合洞察。",
    ], size=15)
    add_panel(slide, 17.0, 2.3, 15.0, 7.2, "项目直观预览")
    add_picture(slide, ASSETS / "02 数据集" / "01 数据集截图.png", 17.5, 3.15, 14.1, 5.8)
    add_callout(slide, 1.0, 10.2, 10.1, 3.8, "研究对象特征", "餐饮平台数据同时具备评分、位置、菜系和服务描述，是少见的多模态城市生活数据。")
    add_callout(slide, 11.5, 10.2, 10.1, 3.8, "研究问题", "哪些区域更成熟？哪些业态更受欢迎？哪些特征能预测更高评分？")
    add_callout(slide, 22.0, 10.2, 10.0, 3.8, "答辩价值", "既展示工程实现，也展示数据分析和业务解释能力。")
    write_notes(slide, "研究背景页建议从平台数据价值切入，再落到 Chennai 的城市特征，最后引出本课题的研究问题。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "01", "研究意义与项目目标", "明确课题对消费者、商家和数据分析场景的价值。", start_page + 1)
    add_panel(slide, 1.0, 2.4, 10.2, 10.9, "研究意义")
    add_bullets(slide, 1.35, 3.35, 9.4, 8.9, [
        "市场洞察：帮助餐饮从业者理解不同市场定位对评分的影响。",
        "地理分析：揭示 Chennai 各区域的餐饮供给密度与评分空间异质性。",
        "消费者决策：通过可视化方式提供区域、菜系和服务的直观比较。",
        "商业选址：利用供给量与平均评分的双轴关系辅助发现潜力区。",
    ], size=14)
    add_panel(slide, 11.8, 2.4, 10.2, 10.9, "项目目标")
    add_bullets(slide, 12.15, 3.35, 9.4, 8.9, [
        "全面描述 Chennai 餐饮市场的整体格局。",
        "分析 market_segment、area、cuisine 与 rating 的关联关系。",
        "识别高供给高评分、低供给高评分等战略性区域。",
        "构建机器学习模型预测评分并解释关键驱动因素。",
    ], size=14)
    add_panel(slide, 22.6, 2.4, 9.6, 10.9, "项目成果形式")
    add_bullets(slide, 22.95, 3.35, 8.8, 8.9, [
        "Flask 交互式 Web 可视化平台。",
        "覆盖描述性统计、空间分析、关联分析、评分预测四个层次。",
        "多图表、可筛选、可解释的城市餐饮分析样例。",
        "为后续推荐系统与选址模型提供基础。"
    ], size=14)
    write_notes(slide, "这一页可以按研究意义、研究目标、成果形式三列讲解，结尾强调项目是分析平台而不是单一页面展示。")


def dataset_slides(prs, start_page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "02", "数据来源与规模", "先说明研究证据来自哪里，再说明它的覆盖范围。", start_page)
    add_panel(slide, 1.0, 2.5, 11.0, 4.0, "数据来源")
    add_bullets(slide, 1.3, 3.3, 10.3, 2.6, [
        "Kaggle 数据集：devvraj/chennai-restaurant-dataset。",
        "原始数据采自 Zomato 平台的 Chennai 餐厅公开信息。",
        "经过清洗后形成适合分析与建模的主表。"
    ], size=14)
    add_panel(slide, 13.0, 2.5, 19.0, 4.0, "核心规模")
    add_kpi(slide, 13.4, 3.3, 4.1, 2.4, "总餐厅数", "11,848", "主表样本量")
    add_kpi(slide, 17.95, 3.3, 4.1, 2.4, "覆盖区域", "268", "可做空间比较")
    add_kpi(slide, 22.5, 3.3, 4.1, 2.4, "市场细分", "8 类", "业态分析基础")
    add_kpi(slide, 27.05, 3.3, 4.1, 2.4, "平均评分", "3.52", "平台评分中枢")
    add_panel(slide, 1.0, 7.1, 15.0, 8.0, "数据与平台快照")
    add_picture(slide, ASSETS / "02 数据集" / "01 数据集截图.png", 1.4, 8.0, 14.2, 6.4)
    add_panel(slide, 17.0, 7.1, 15.0, 8.0, "数据质量速览")
    add_picture(slide, ASSETS / "02 数据集" / "01_scorecard.png", 17.4, 8.0, 14.2, 6.4)
    write_notes(slide, "数据来源与规模页主要建立可信度，尤其要强调样本量、区域覆盖和市场细分数，为后续所有图表做铺垫。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "02", "核心字段与数据结构", "解释后续图表和模型为什么能成立。", start_page + 1)
    add_panel(slide, 1.0, 2.4, 15.3, 12.5, "核心字段")
    field_rows = [
        ("restaurant", "文本", "餐厅名称"),
        ("rating", "数值", "目标变量，用户评分"),
        ("market_segment", "分类", "业态细分，如 Dine-out / Cafes / Bar"),
        ("area", "分类", "所在区域，如 Porur / Anna Nagar / T. Nagar"),
        ("cuisine", "文本", "菜系列表"),
        ("top_dishes", "文本", "招牌菜列表"),
        ("features", "文本", "服务特色"),
        ("latitude / longitude", "数值", "空间坐标"),
    ]
    y = 3.2
    for name, typ, desc in field_rows:
        add_rect(slide, 1.35, y, 14.6, 1.0, fill=WHITE, line=LINE, radius=False)
        add_textbox(slide, 1.55, y + 0.18, 3.6, 0.5, name, size=12, bold=True, color=BLUE_2, font=FONT_EN if "_" in name or "/" in name else FONT_CN)
        add_textbox(slide, 5.25, y + 0.18, 2.4, 0.5, typ, size=11, color=TEXT)
        add_textbox(slide, 7.55, y + 0.18, 7.9, 0.5, desc, size=11, color=MUTED)
        y += 1.2
    add_panel(slide, 17.0, 2.4, 15.0, 5.2, "为什么这些字段重要")
    add_bullets(slide, 17.35, 3.25, 14.2, 3.4, [
        "rating 是分析与预测的核心目标变量。",
        "area 和经纬度提供空间分析能力。",
        "market_segment 决定业态差异比较的主轴。",
        "cuisine、top_dishes、features 提供文本语义和服务供给信息。",
    ], size=13)
    add_panel(slide, 17.0, 8.1, 15.0, 6.8, "数据结构理解")
    add_bullets(slide, 17.35, 8.95, 14.1, 5.0, [
        "主表适合做总体统计、地图散点和区域聚合。",
        "多值文本字段无法直接做频率和交叉比较，因此必须拆成长格式表。",
        "这也是后续能做热力图、菜系排行、招牌菜分析和文本建模的基础。"
    ], size=13)
    write_notes(slide, "这一页要说明字段不是简单列举，而是直接对应到后续分析方法。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "02", "数据预处理流程", "把原始平台字段转化为可分析、可建模的结构化证据。", start_page + 2)
    add_panel(slide, 1.0, 2.6, 10.6, 12.0, "处理步骤")
    steps = [
        "检查数据完整性与重复值",
        "列名规范化：统一小写、下划线风格",
        "验证字段范围与坐标有效性",
        "构造 restaurant_id、rating_band 等衍生列",
        "为文本与空间分析准备长格式数据和空间特征",
    ]
    y = 3.45
    for i, step in enumerate(steps, start=1):
        add_chip(slide, 1.45, y - 0.05, 1.0, 0.7, str(i).zfill(2), fill=BLUE)
        add_rect(slide, 2.7, y - 0.08, 8.3, 0.82, fill=WHITE, line=LINE)
        add_textbox(slide, 3.0, y + 0.1, 7.6, 0.4, step, size=12, color=TEXT)
        y += 1.75
    add_panel(slide, 12.1, 2.6, 9.7, 12.0, "关键质量指标")
    add_bullets(slide, 12.45, 3.45, 8.8, 10.6, [
        "主数据集 11,848 行、11 列，无缺失评分。",
        "评分最小值 0.3，最大值 4.9，中位数 3.6。",
        "唯一餐厅名称 8,352 个，说明存在较多连锁或同名门店。",
        "覆盖 268 个 area，空间分布粒度足够细。",
        "坐标范围完整，适合地图和空间热力分析。",
    ], size=13)
    add_panel(slide, 22.3, 2.6, 9.7, 12.0, "评分分段设计")
    add_bullets(slide, 22.65, 3.45, 8.8, 10.6, [
        "fragile: rating < 2.9",
        "developing: 2.9 ≤ rating < 3.4",
        "solid: 3.4 ≤ rating < 3.8",
        "strong: 3.8 ≤ rating < 4.2",
        "elite: rating ≥ 4.2",
        "该分段既能支撑分组柱形图，也方便后续区域堆叠分析。",
    ], size=13)
    write_notes(slide, "预处理流程页建议强调为什么要设计 rating_band，它让不同图表之间的度量体系保持一致。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "02", "长格式拆分与特征工程", "解决多值文本字段无法直接分析的问题。", start_page + 3)
    add_panel(slide, 1.0, 2.5, 16.0, 12.0, "长格式数据拆分")
    add_bullets(slide, 1.4, 3.35, 15.0, 3.2, [
        "原始数据中的 cuisine、features、top_dishes 都是逗号分隔的多值字段，无法直接用于频率统计和交叉分析。",
        "因此把它们展开为 cuisine_long、feature_long、dish_long 三张长格式表。"
    ], size=13)
    add_rect(slide, 1.4, 6.3, 15.1, 6.0, fill=WHITE, line=LINE)
    headers = [("维度", 1.7), ("长格式表", 5.2), ("行数", 10.0), ("唯一值", 12.0), ("膨胀倍数", 14.4)]
    for text, x in headers:
        add_textbox(slide, x, 6.55, 2.1, 0.45, text, size=11, bold=True, color=BLUE)
    rows = [
        ("菜系", "long_format_cuisine.csv", "28,799", "92", "2.43x"),
        ("特色", "long_format_feature.csv", "32,698", "78", "2.76x"),
        ("菜品", "long_format_dish.csv", "76,763", "2,036", "6.48x"),
    ]
    y = 7.45
    for r in rows:
        add_rect(slide, 1.6, y, 14.6, 1.2, fill=PALE_BG, line=LINE)
        vals = [r[0], r[1], r[2], r[3], r[4]]
        xs = [1.8, 5.0, 10.0, 12.1, 14.55]
        ws = [2.0, 4.6, 1.5, 1.2, 1.3]
        for val, x, w in zip(vals, xs, ws):
            add_textbox(slide, x, y + 0.27, w, 0.5, val, size=10, color=TEXT, bold=x == 1.8)
        y += 1.45
    add_panel(slide, 17.5, 2.5, 14.5, 12.0, "建模特征工程")
    add_bullets(slide, 17.9, 3.35, 13.7, 10.0, [
        "空间偏移：lat_offset、lon_offset，刻画相对城市中心的地理位置。",
        "距离特征：distance_to_city_median_km，量化离城市中心的远近。",
        "品牌足迹：same_name_outlets、is_multi_outlet_name，识别连锁品牌效应。",
        "文本向量化：CountVectorizer / TfidfVectorizer 提取 cuisine、features、top_dishes 语义。",
        "分类编码：market_segment / area 使用 One-Hot 编码，进入机器学习流水线。",
    ], size=13)
    write_notes(slide, "这一页的重点是解释为什么长格式拆分和特征工程是后续可视化与机器学习的共同基础。")


def tech_slides(prs, start_page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "03", "技术栈总览", "采用轻量但完整的前后端与分析技术组合。", start_page)
    add_panel(slide, 1.0, 2.6, 7.6, 5.0, "后端框架")
    add_bullets(slide, 1.3, 3.35, 6.9, 3.3, [
        "Flask：轻量级 Python Web 框架。",
        "Blueprint：组织多页面路由。",
        "分层设计：Service → Route → Template。",
    ], size=13)
    add_panel(slide, 9.2, 2.6, 7.6, 5.0, "数据处理")
    add_bullets(slide, 9.5, 3.35, 6.9, 3.3, [
        "Pandas：清洗、聚合与表结构变换。",
        "NumPy：基础统计与直方图计算。",
        "SciPy：核密度估计（KDE）。",
    ], size=13)
    add_panel(slide, 17.4, 2.6, 7.6, 5.0, "可视化")
    add_bullets(slide, 17.7, 3.35, 6.9, 3.3, [
        "Plotly.js：交互式图表主力。",
        "ECharts：词云图。",
        "Mapbox：地理空间可视化。",
    ], size=13)
    add_panel(slide, 25.6, 2.6, 6.4, 5.0, "机器学习")
    add_bullets(slide, 25.9, 3.35, 5.7, 3.3, [
        "scikit-learn",
        "Ridge / ElasticNet",
        "RandomForest / HGB",
    ], size=13)
    add_panel(slide, 1.0, 8.2, 15.4, 6.2, "前端与交互")
    add_bullets(slide, 1.3, 9.0, 14.7, 4.8, [
        "原生 JavaScript，无前端框架依赖，降低部署复杂度。",
        "Plotly.js 通过 CDN 引入，支持地图、柱形图、热力图、箱线图、六边形密度图等。",
        "自定义 CSS 主题为 Chennai 餐饮数据构建统一暖色视觉风格。",
    ], size=13)
    add_panel(slide, 17.0, 8.2, 15.0, 6.2, "建模工程")
    add_bullets(slide, 17.3, 9.0, 14.1, 4.8, [
        "Pipeline + ColumnTransformer：统一预处理与训练流程。",
        "GroupShuffleSplit：按餐厅名称分组划分数据，避免数据泄漏。",
        "Permutation Importance、PDP 与 Ridge 系数：支撑模型可解释性分析。",
    ], size=13)
    write_notes(slide, "技术栈页不需要逐一介绍库，而是要说明这些工具如何共同支撑多页面分析与评分预测。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "03", "系统架构与页面组织", "系统结构对应答辩中的分析流程。", start_page + 1)
    add_panel(slide, 1.0, 2.4, 18.0, 12.0, "系统架构")
    layers = [
        ("用户浏览器", "浏览与交互分析页面"),
        ("Flask (run.py)", "统一入口与页面调度"),
        ("Route 层", "overview / spatial / segment / prediction / dashboard"),
        ("Service 层", "OverviewService / SpatialService / SegmentService / PredictionService"),
        ("数据层", "主 CSV + 长格式表 + 缓存结果"),
    ]
    y = 3.2
    for i, (name, desc) in enumerate(layers):
        add_rect(slide, 5.0, y, 10.0, 1.35, fill=WHITE, line=BLUE if i % 2 == 0 else BLUE_2, radius=True)
        add_textbox(slide, 5.25, y + 0.18, 3.4, 0.45, name, size=13, bold=True, color=BLUE)
        add_textbox(slide, 8.7, y + 0.18, 5.9, 0.45, desc, size=11, color=MUTED)
        if i < len(layers) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, cm(9.25), cm(y + 1.4), cm(1.35), cm(0.9))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GOLD
            arrow.line.fill.background()
        y += 2.15
    add_panel(slide, 20.1, 2.4, 11.9, 12.0, "页面与分析层次")
    add_bullets(slide, 20.45, 3.25, 11.1, 10.2, [
        "Overview：整体市场格局与基础评分分布。",
        "Spatial Diagnostics：区域供给、热力图、四象限。",
        "Market Segment Rating：业态评分分布、服务采用率、菜品和菜系结构。",
        "Prediction Analysis：模型排行榜、残差、误差分解、特征重要性。",
        "Dashboard：面向展示的一体化入口页。",
    ], size=13)
    write_notes(slide, "系统架构页可以帮助答辩老师理解项目不是静态报告，而是由页面、服务层和数据层组成的完整系统。")


def eda_slides(prs, start_page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "04", "整体概况：市场全貌", "先从全局地图和评分分布建立市场直觉。", start_page)
    add_panel(slide, 1.0, 2.4, 20.6, 12.0, "Overview 页面")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "overview" / "overview.png", 1.4, 3.2, 19.8, 10.4)
    add_panel(slide, 22.1, 2.4, 10.0, 12.0, "关键信息")
    add_bullets(slide, 22.45, 3.25, 9.2, 8.2, [
        "地图揭示餐厅主要集中在 Chennai 核心城区与沿海带。",
        "平均评分约 3.52，整体偏高，但仍存在显著结构差异。",
        "细分市场和区域直方图为后续深入分析确定了重点对象。",
    ], size=13)
    add_callout(slide, 22.45, 11.9, 8.8, 1.9, "答辩提醒", "这页重点是“建立全局认知”，不要把细节讲太满。", fill=SOFT_RED)
    write_notes(slide, "整体概况页主要回答两个问题：市场供给在哪里，以及评分大致落在哪个区间。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "04", "整体概况：评分分布拆解", "把总览中的关键图单独放大，解决上一版图太小的问题。", start_page + 1)
    add_panel(slide, 1.0, 2.4, 15.5, 6.1, "Market Segment 评分直方图")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "overview" / "Market Segment 评分直方图.png", 1.4, 3.2, 14.7, 4.7)
    add_panel(slide, 17.0, 2.4, 15.0, 6.1, "Area 评分直方图")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "overview" / "Area 评分直方图.png", 17.4, 3.2, 14.2, 4.7)
    add_panel(slide, 1.0, 9.0, 15.5, 5.4, "Area 评分直方图补充")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "overview" / "Area 评分直方图2.png", 1.4, 9.8, 14.7, 4.0)
    add_panel(slide, 17.0, 9.0, 15.0, 5.4, "读图结论")
    add_bullets(slide, 17.35, 9.85, 14.1, 3.8, [
        "Cloud Kitchen 在较高评分段表现更优。",
        "部分高密度区域的评分分布更分散，说明竞争更激烈。",
        "把这些图拆页展示后，更适合现场指图说明差异。"
    ], size=13)
    write_notes(slide, "这一页的目的就是把图放大讲清楚，避免上一版所有图堆在一页看不清。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "04", "空间分析：区域分布与评分结构", "从区域名字、评分分布到等级构成，逐层理解空间异质性。", start_page + 2)
    add_panel(slide, 1.0, 2.4, 10.2, 5.9, "Area 词云")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "空间差异性分析" / "词云图.png", 1.35, 3.1, 9.5, 4.6)
    add_panel(slide, 11.8, 2.4, 10.2, 5.9, "区域评分箱线图")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "空间差异性分析" / "箱线图.png", 12.15, 3.1, 9.5, 4.6)
    add_panel(slide, 22.6, 2.4, 9.4, 5.9, "评分等级堆叠图")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "空间差异性分析" / "堆叠图.png", 22.95, 3.1, 8.7, 4.6)
    add_panel(slide, 1.0, 9.0, 31.0, 5.4, "空间读图结论")
    add_bullets(slide, 1.35, 9.85, 30.1, 3.8, [
        "词云图直观显示餐厅更集中的区域名称，便于先从“供给量”建立直觉。",
        "箱线图展示出不同 area 的评分中位数、波动范围和异常值差异。",
        "堆叠图把 fragile 到 elite 的分段结构进一步展开，体现区域内部的质量构成差异。"
    ], size=13)
    write_notes(slide, "空间分析第一页建议按‘区域名字出现频率 → 区域评分分布 → 区域等级结构’的顺序讲。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "04", "空间分析：供给、热力与战略分区", "这是项目中最具决策意味的一组空间证据。", start_page + 3)
    add_panel(slide, 1.0, 2.4, 10.2, 5.8, "区域供给 vs 平均评分 气泡图")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "空间差异性分析" / "气泡图.png", 1.35, 3.1, 9.5, 4.5)
    add_panel(slide, 11.8, 2.4, 20.2, 5.8, "核函数热力图")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "空间差异性分析" / "核函数热力图.png", 12.15, 3.1, 19.5, 4.5)
    add_panel(slide, 1.0, 8.8, 19.0, 5.6, "四象限矩阵")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "空间差异性分析" / "四象限图.png", 1.35, 9.55, 18.3, 4.2)
    add_panel(slide, 20.4, 8.8, 11.6, 5.6, "战略结论")
    add_bullets(slide, 20.75, 9.55, 10.8, 4.2, [
        "T. Nagar / Anna Nagar 属于成熟优质区。",
        "Porur 等高供给低评分区域竞争更激烈。",
        "低供给高评分区域可能存在商业机会。",
    ], size=12)
    write_notes(slide, "这里要突出供给量和平均评分不是简单正相关，四象限因此具备了选址与市场判断的实际意义。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "04", "业态分析：评分分布与数量结构", "从 market segment 角度看不同业态的评分表现。", start_page + 4)
    add_panel(slide, 1.0, 2.4, 15.3, 6.0, "评分脊线图")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "segment对评分影响" / "脊图.png", 1.35, 3.2, 14.6, 4.6)
    add_panel(slide, 17.0, 2.4, 15.0, 6.0, "评分等级分组柱形图")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "segment对评分影响" / "分组柱形图.png", 17.35, 3.2, 14.2, 4.6)
    add_panel(slide, 1.0, 8.9, 31.0, 5.6, "读图结论")
    add_bullets(slide, 1.35, 9.75, 30.2, 3.8, [
        "脊线图更适合比较各业态评分分布的形状和中位数位置。",
        "分组柱形图则把数量差异讲清楚，尤其能看出 Restaurant / Fast Food 的规模优势。",
        "Cloud Kitchen 在高评分段的占比更突出，是后续解释的关键发现。",
    ], size=13)
    write_notes(slide, "这一页要强调‘分布形态’和‘数量结构’是两件不同的事，所以拆成两张大图会更清楚。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "04", "业态分析：服务特征、菜系与招牌菜", "把用户感知层面的信息补上。", start_page + 5)
    add_panel(slide, 1.0, 2.4, 14.8, 5.9, "服务特色采用率热力图")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "segment对评分影响" / "菜系结构提供服务特色比例图.png", 1.35, 3.15, 14.1, 4.7)
    add_panel(slide, 16.3, 2.4, 15.7, 5.9, "菜系与招牌菜结构")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "segment对评分影响" / "最受喜爱的菜品和当地菜品图.png", 16.7, 3.15, 14.9, 4.7)
    add_panel(slide, 1.0, 8.9, 31.0, 5.6, "结论与解释")
    add_bullets(slide, 1.35, 9.75, 30.1, 3.9, [
        "Home Delivery 和 Indoor Seating 是最普遍的两类服务特征，分别对应便利性与到店体验。",
        "Chinese 菜系最常见，而 Dosa、Idli、Filter Coffee 等本地菜在评分和覆盖率上同样表现突出。",
        "这说明评分不仅受区域影响，也受业态组织方式和本地饮食文化偏好的共同作用。"
    ], size=13)
    write_notes(slide, "这一页负责把服务特征、菜系和招牌菜串起来，突出‘文化偏好 + 服务模式’对评分的影响。")


def model_slides(prs, start_page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "05", "评分预测：模型比较", "先交代模型路线，再说明最优模型表现。", start_page)
    add_panel(slide, 1.0, 2.5, 12.5, 4.1, "候选模型")
    add_bullets(slide, 1.35, 3.3, 11.8, 2.5, [
        "DummyRegressor（基线）",
        "Ridge / ElasticNet（线性基准）",
        "RandomForest（树模型）",
        "HistGradientBoosting（最终表现最优）",
    ], size=14)
    add_panel(slide, 14.2, 2.5, 17.8, 11.8, "Model Leaderboard")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "predict" / "模型排行榜.png", 14.6, 3.35, 17.0, 10.3)
    add_panel(slide, 1.0, 7.1, 12.5, 7.2, "解释重点")
    add_bullets(slide, 1.35, 7.95, 11.7, 5.4, [
        "评价指标以 MAE 为主，同时观察 Train MAE 与 Test MAE。",
        "梯度提升模型在泛化误差上更优，说明非线性关系较强。",
        "答辩时不必强调模型多复杂，而要强调模型确实比线性基准捕捉到了更多结构。",
    ], size=13)
    write_notes(slide, "模型比较页要把重点放在‘为什么需要更强模型’和‘最优模型相比基线的提升’。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "05", "评分预测：残差与误差分解", "看模型在哪里预测得好，哪里还有系统性误差。", start_page + 1)
    add_panel(slide, 1.0, 2.4, 10.2, 5.8, "六边形密度图")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "predict" / "六边形密度图.png", 1.35, 3.15, 9.5, 4.5)
    add_panel(slide, 11.8, 2.4, 10.2, 5.8, "残差分布")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "predict" / "残差分析.png", 12.15, 3.15, 9.5, 4.5)
    add_panel(slide, 22.6, 2.4, 9.4, 5.8, "Area 分解 MAE + Bias")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "predict" / "Area 分解 MAE + Bias.png", 22.95, 3.15, 8.7, 4.5)
    add_panel(slide, 1.0, 8.8, 14.8, 5.6, "Market Segment 分解 MAE + Bias")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "predict" / "Market Segment分解 MAE + Bias.png", 1.35, 9.55, 14.1, 4.2)
    add_panel(slide, 16.3, 8.8, 15.7, 5.6, "误差解释")
    add_bullets(slide, 16.65, 9.55, 14.8, 4.1, [
        "预测点主要集中在 3.4~3.8 评分区间，说明模型对主流样本学习较稳定。",
        "残差整体围绕 0 分布，但极端样本依旧更难拟合。",
        "不同区域和市场细分的 MAE / Bias 不同，说明评分机制存在分层差异。"
    ], size=12)
    write_notes(slide, "这一页讲‘模型误差不是随机噪声’，而是和区域、业态结构有关。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "05", "评分预测：特征重要性与机制解释", "把机器学习结果转回可理解的业务语言。", start_page + 2)
    add_panel(slide, 1.0, 2.4, 15.3, 6.0, "Permutation Importance")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "predict" / "重要性图.png", 1.35, 3.15, 14.6, 4.7)
    add_panel(slide, 17.0, 2.4, 15.0, 6.0, "Ridge 系数方向")
    add_picture(slide, ASSETS / "04 EDA数据分析" / "predict" / "脊线系数方向.png", 17.35, 3.15, 14.2, 4.7)
    add_panel(slide, 1.0, 8.9, 31.0, 5.6, "核心发现")
    add_bullets(slide, 1.35, 9.75, 30.2, 3.9, [
        "features_tokens、top_dishes_tokens 和 same_name_outlets 是最关键的特征块。",
        "品牌足迹显著影响评分，说明连锁品牌更容易建立稳定口碑。",
        "菜单语义和招牌菜丰富度同样重要，证明“提供什么”和“如何表达”都会影响用户评分感知。",
        "距离市中心更远的餐厅评分略低，空间位置仍是背景变量。"
    ], size=13)
    write_notes(slide, "最后一页模型解释要把技术指标翻译成业务语言，答辩老师更关心你从模型里读出了什么。")


def summary_slides(prs, start_page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "06", "主要结论", "把前面的图表与模型结果收束成清晰论点。", start_page)
    add_panel(slide, 1.0, 2.6, 10.0, 11.6, "结论 1")
    add_textbox(slide, 1.35, 3.3, 9.1, 0.7, "Chennai 餐饮市场以 Dine-out 为主导，Cloud Kitchen 评分更高。", size=16, bold=True, color=TEXT)
    add_bullets(slide, 1.35, 4.4, 9.0, 7.0, [
        "总体规模上，Restaurant / Fast Food 数量最大。",
        "但在评分分布上，Cloud Kitchen 更容易进入较高评分段。"
    ], size=13)
    add_panel(slide, 11.7, 2.6, 10.0, 11.6, "结论 2")
    add_textbox(slide, 12.05, 3.3, 9.1, 0.7, "空间分布不均匀，T. Nagar / Anna Nagar 是核心优质区。", size=16, bold=True, color=TEXT)
    add_bullets(slide, 12.05, 4.4, 9.0, 7.0, [
        "高供给高评分区域已经形成成熟市场。",
        "部分高供给低评分区域竞争更激烈，选址风险更高。"
    ], size=13)
    add_panel(slide, 22.4, 2.6, 9.6, 11.6, "结论 3")
    add_textbox(slide, 22.75, 3.3, 8.8, 0.7, "品牌连锁效应对评分有显著正向影响。", size=16, bold=True, color=TEXT)
    add_bullets(slide, 22.75, 4.4, 8.4, 7.0, [
        "同名门店数是模型最强正向因子之一。",
        "机器学习结果验证了品牌足迹的重要性。"
    ], size=13)
    write_notes(slide, "主要结论页建议一条结论配一层证据回顾，让老师能快速回想前面哪几张图支撑了这条话。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_header(slide, "06", "局限性与展望", "最后用一页完成研究反思和后续延伸。", start_page + 1)
    add_panel(slide, 1.0, 2.6, 15.2, 11.8, "局限性")
    add_bullets(slide, 1.35, 3.45, 14.4, 9.8, [
        "数据时效性有限，平台评分会随时间变化。",
        "尚未纳入价格区间、评论数、图片数量等更丰富的行为变量。",
        "当前分析以截面数据为主，未深入展开时间序列变化。",
        "模型解释已较充分，但仍可结合更细粒度文本语义进一步提升表现。"
    ], size=14)
    add_panel(slide, 17.0, 2.6, 15.0, 11.8, "后续工作")
    add_bullets(slide, 17.35, 3.45, 14.2, 9.8, [
        "增加价格、评论量、照片量等变量，提升分析深度。",
        "加入时间维度，追踪评分和供给变化趋势。",
        "延伸到推荐系统，实现基于位置与菜系偏好的个性化推荐。",
        "尝试将选址判断与商圈分析结合，形成更直接的经营决策工具。"
    ], size=14)
    write_notes(slide, "最后一页保持谦逊和完整性，既说明现阶段成果，也自然引出后续扩展方向。")


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = cm(33.867)
    prs.slide_height = cm(19.05)

    cover_slide(prs, 1)
    agenda_slide(prs, 2)
    background_slides(prs, 3)
    dataset_slides(prs, 5)
    tech_slides(prs, 9)
    eda_slides(prs, 11)
    model_slides(prs, 17)
    summary_slides(prs, 20)
    return prs


def main():
    ensure_output()
    prs = build_presentation()
    prs.save(str(PPT_PATH))
    print(PPT_PATH)


if __name__ == "__main__":
    main()
